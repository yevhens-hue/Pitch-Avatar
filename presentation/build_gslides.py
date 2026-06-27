#!/usr/bin/env python3
"""
Builds an editable, Google-Slides-compatible .pptx (16:9) from the improved
competitor review content. Native text boxes / shapes / table (not images),
so everything stays editable after importing into Google Slides.
Run:  python3 presentation/build_gslides.py
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

HERE = os.path.dirname(os.path.abspath(__file__))
ASSETS = os.path.join(HERE, "assets")
OUT = os.path.join(HERE, "competitor_review_v3.pptx")

# ---- palette ----
BRAND   = RGBColor(0x0B, 0x5F, 0xFF)
INK     = RGBColor(0x0F, 0x17, 0x2A)
MUTED   = RGBColor(0x64, 0x74, 0x8B)
LINE    = RGBColor(0xE5, 0xE9, 0xF0)
SOFT    = RGBColor(0xEE, 0xF4, 0xFF)
SOFTBD  = RGBColor(0xD6, 0xE4, 0xFF)
ALT     = RGBColor(0xF7, 0xF9, 0xFC)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
GREEN   = RGBColor(0x16, 0xA3, 0x4A)
AMBER   = RGBColor(0xD9, 0x77, 0x06)
AMBERBD = RGBColor(0xF6, 0xD9, 0xA8)
VIOLET  = RGBColor(0x7C, 0x3A, 0xED)
COVER1  = RGBColor(0x0B, 0x2E, 0x8A)
COVER2  = RGBColor(0x07, 0x1A, 0x52)
INK2    = RGBColor(0x33, 0x41, 0x55)
NAVYTXT = RGBColor(0x1E, 0x3A, 0x8A)
BLUETXT = RGBColor(0xC9, 0xD8, 0xFF)
KICK_ON = RGBColor(0x9F, 0xC0, 0xFF)

FONT = "Inter"
EMU_IN = 914400

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = 13.333, 7.5


def slide():
    return prs.slides.add_slide(BLANK)


def _set_fill(shape, color):
    if color is None:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = color


def _set_line(shape, color, w=1.0):
    if color is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = color
        shape.line.width = Pt(w)


def rect(s, x, y, w, h, fill, line=None, lw=1.0, rounded=False, radius=0.08, shadow=False):
    shp = s.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h))
    _set_fill(shp, fill)
    _set_line(shp, line, lw)
    shp.shadow.inherit = False
    if rounded:
        try:
            shp.adjustments[0] = radius
        except Exception:
            pass
    return shp


def text(s, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
         space_after=4, line_spacing=1.0):
    """runs: list of paragraphs; each paragraph is list of (txt,size,color,bold,italic,tracking)."""
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space_after)
        p.space_before = Pt(0)
        p.line_spacing = line_spacing
        for (t, size, color, bold, *rest) in para:
            italic = rest[0] if len(rest) > 0 else False
            r = p.add_run()
            r.text = t
            r.font.name = FONT
            r.font.size = Pt(size)
            r.font.bold = bold
            r.font.italic = italic
            r.font.color.rgb = color
    return tb


def kicker(s, label, num):
    dot = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.92), Inches(0.62), Inches(0.09), Inches(0.09))
    _set_fill(dot, BRAND); _set_line(dot, None); dot.shadow.inherit = False
    text(s, 1.12, 0.5, 9.5, 0.4, [[(label.upper(), 11.5, BRAND, True)]])
    text(s, SW - 1.6, 0.5, 0.7, 0.4, [[(num, 11.5, MUTED, True)]], align=PP_ALIGN.RIGHT)


def title(s, t, y=0.95):
    text(s, 0.92, y, 11.5, 0.8, [[(t, 27, INK, True)]])


def lead(s, t, y=1.62, w=10.6):
    text(s, 0.92, y, w, 0.8, [[(t, 13.5, MUTED, False)]], line_spacing=1.15)


def footer(s, right="Обзор конкурентов · v3"):
    rect(s, 0.92, 6.95, SW - 1.84, 0.012, LINE, None)
    play(s, 0.92, 7.06, 0.16)
    text(s, 1.14, 7.04, 4, 0.3, [[("Pitch Avatar", 10.5, INK, True)]])
    text(s, SW - 5.92, 7.04, 5, 0.3, [[(right, 9.5, RGBColor(0x93,0xA0,0xB4), False)]], align=PP_ALIGN.RIGHT)


def play(s, x, y, size, bg=BRAND, tri=WHITE):
    sq = rect(s, x, y, size, size, bg, None, rounded=True, radius=0.3)
    tr = s.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, Inches(x+size*0.32), Inches(y+size*0.26),
                            Inches(size*0.38), Inches(size*0.48))
    _set_fill(tr, tri); _set_line(tr, None); tr.shadow.inherit = False
    tr.rotation = 90
    return sq


def card(s, x, y, w, h, tag=None, tag_color=BRAND, tag_bg=SOFT, heading=None, body=None,
         bullets=None, fill=WHITE, border=LINE):
    rect(s, x, y, w, h, fill, border, 1.0, rounded=True, radius=0.06)
    cy = y + 0.28
    if tag:
        tw = 0.18 + len(tag) * 0.085
        rect(s, x + 0.28, cy, tw, 0.32, tag_bg, None, rounded=True, radius=0.5)
        text(s, x + 0.28, cy + 0.02, tw, 0.3, [[(tag, 10, tag_color, True)]], align=PP_ALIGN.CENTER)
        cy += 0.52
    if heading:
        text(s, x + 0.28, cy, w - 0.56, 0.4, [[(heading, 16, INK, True)]])
        cy += 0.46
    if body:
        text(s, x + 0.28, cy, w - 0.56, h - (cy - y) - 0.2,
             [[(body, 12, MUTED, False)]], line_spacing=1.18)
    if bullets:
        paras = []
        for b in bullets:
            paras.append([("•  ", 12, BRAND, True), (b, 12, INK2, False)])
        text(s, x + 0.28, cy, w - 0.56, h - (cy - y) - 0.2, paras, space_after=7, line_spacing=1.1)


def takeaway(s, x, y, w, label, body, h=0.95):
    rect(s, x, y, w, h, SOFT, SOFTBD, 1.0, rounded=True, radius=0.12)
    text(s, x + 0.32, y + 0.18, w - 0.64, h - 0.36,
         [[(label + " ", 13, BRAND, True), (body, 13, NAVYTXT, False)]],
         anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.2)


def annot(s, x, y, w, num, head, body):
    b = rect(s, x, y, 0.34, 0.34, INK, None, rounded=True, radius=0.28)
    text(s, x, y + 0.02, 0.34, 0.32, [[(num, 13, WHITE, True)]], align=PP_ALIGN.CENTER,
         anchor=MSO_ANCHOR.MIDDLE)
    text(s, x + 0.5, y - 0.04, w - 0.5, 0.35, [[(head, 14, INK, True)]])
    text(s, x + 0.5, y + 0.32, w - 0.5, 0.6, [[(body, 12, MUTED, False)]], line_spacing=1.15)


# =================================================================== 1 COVER
s = slide()
rect(s, 0, 0, SW, SH, COVER1, None)
rect(s, 0, 4.7, SW, SH - 4.7, COVER2, None)
glow = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(9.6), Inches(0.6), Inches(5.2), Inches(5.2))
_set_fill(glow, RGBColor(0x1E, 0x4F, 0xD8)); _set_line(glow, None); glow.shadow.inherit = False
play(s, 0.92, 0.62, 0.34, bg=WHITE, tri=BRAND)
text(s, 1.38, 0.64, 5, 0.4, [[("Pitch Avatar", 18, WHITE, True)]])
text(s, 0.92, 2.0, 11, 0.4, [[("PRODUCT STRATEGY · COMPETITIVE TEARDOWN", 12.5, KICK_ON, True)]])
text(s, 0.9, 2.45, 11.6, 2.2, [
    [("Use Cases & UI Skins:", 44, WHITE, True)],
    [("как продавать сложный продукт", 44, WHITE, True)],
    [("разным аудиториям", 44, WHITE, True)],
], line_spacing=1.02)
text(s, 0.92, 4.95, 9, 0.8, [[("Что делают Monday, ClickUp, HubSpot и Salesforce — и какие из этих практик стоит перенести в Pitch Avatar.", 15, BLUETXT, False)]], line_spacing=1.3)
meta = [("Документ", "Обзор конкурентов v3"), ("Фокус", "Onboarding · Modular UI · Admin"),
        ("Аудитории", "HR · Marketing · Sales")]
mx = 0.92
for k, v in meta:
    text(s, mx, 5.95, 3.6, 0.3, [[(k, 11, RGBColor(0xAF,0xC4,0xFF), False)]])
    text(s, mx, 6.25, 3.6, 0.4, [[(v, 13, WHITE, True)]])
    mx += 3.9

# =================================================================== 2 EXEC SUMMARY
s = slide()
kicker(s, "Executive Summary", "02")
title(s, "Коротко: один вопрос на входе + модульный интерфейс")
lead(s, "Лидеры рынка решают конфликт «мощность против простоты» двумя приёмами, которые работают в связке. Это основа рекомендаций для Pitch Avatar.")
card(s, 0.92, 2.35, 5.55, 1.85, tag="Паттерн 1 · Intent", heading="Intent-driven Onboarding",
     body="Спросить о цели один раз при регистрации и записать Use Case в профиль. Дальше продукт сам подстраивает интерфейс.", fill=ALT)
card(s, 6.78, 2.35, 5.55, 1.85, tag="Паттерн 2 · Modular", tag_color=VIOLET,
     tag_bg=RGBColor(0xF1,0xEC,0xFE), heading="Role-Based / Modular UI",
     body="Единое ядро данных, поверх которого включаются и выключаются модули. Лишнее скрыто; сложное — за «замочком».", fill=ALT)
stats = [("1", "вопрос на онбординге определяет всю дальнейшую конфигурацию"),
         ("3", "аудитории (HR · Marketing · Sales) — три UI-скина на одном ядре"),
         ("403", "скрытие в UI обязательно подкрепляется защитой бэкенда")]
sx = 0.92
for v, k in stats:
    text(s, sx, 4.55, 3.7, 0.7, [[(v, 36, BRAND, True)]])
    text(s, sx, 5.35, 3.5, 0.9, [[(k, 12, MUTED, False)]], line_spacing=1.15)
    sx += 4.05
footer(s)

# =================================================================== 3 PROBLEM
s = slide()
kicker(s, "Проблема", "03")
title(s, "Complexity vs Usability")
lead(s, "Pitch Avatar — мощный продукт. Но один и тот же интерфейс пугает новичка и тормозит эксперта. Как продать ценность каждой аудитории, не утопив её в функциях?")
card(s, 0.92, 2.45, 5.55, 3.9, tag="Риск", tag_color=AMBER, tag_bg=RGBColor(0xFC,0xF1,0xE0),
     heading="Что ломается сегодня", fill=ALT,
     bullets=["Новичок видит все функции сразу → перегрузка и отток.",
              "HR, маркетолог и продавец получают один и тот же экран.",
              "Ценность продвинутых функций просто не замечают."])
card(s, 6.78, 2.45, 5.55, 3.9, tag="Ответ рынка", heading="Как решают лидеры",
     fill=SOFT, border=SOFTBD,
     bullets=["Intent-driven onboarding — узнать цель и показать релевантное.",
              "Modular UI / role-based workspaces — свой набор модулей.",
              "Progressive disclosure + upsell — сложное «продаётся» тизерами."])
footer(s)

# =================================================================== 4 ONBOARDING
s = slide()
kicker(s, "Кейс 1 · Онбординг", "04")
title(s, "Monday.com и ClickUp: цель спрашивают сразу")
lead(s, "Оба продукта превращают первый экран в инструмент сегментации и используют ответ для маршрутизации по правильной воронке.")
cw = 3.72
cx = 0.92
cards4 = [("Survey на входе", "Обязательный вопрос при регистрации: «Для чего вы будете использовать платформу?». Ответ — в профиль."),
          ("Привязка к доменам", "Кастомные лендинги пробрасывают UTM-параметры и предзаполняют ответы, пропуская лишние шаги."),
          ("Trial vs Demo", "Self-serve trial для малых команд; Enterprise-трафик уходит на Demo с участием sales.")]
for h, b in cards4:
    card(s, cx, 2.4, cw, 2.2, heading=h, body=b)
    cx += cw + 0.27
takeaway(s, 0.92, 4.95, SW - 1.84, "Вывод для Pitch Avatar:",
         "первый экран — точка сегментации. Один вопрос экономит десятки последующих и задаёт стартовую конфигурацию.")
footer(s)

# =================================================================== 5 MOCKUP onboarding
def mockup_slide(num, kick, img, heading, anns, foot):
    s = slide()
    kicker(s, kick, num)
    rect(s, 0.92, 1.25, 6.7, 5.2, ALT, LINE, 1.0, rounded=True, radius=0.04)
    s.shapes.add_picture(os.path.join(ASSETS, img), Inches(1.07), Inches(1.4), width=Inches(6.4))
    text(s, 8.05, 1.5, 4.4, 0.6, [[(heading, 22, INK, True)]], line_spacing=1.1)
    ay = 2.45
    for n, head, body in anns:
        annot(s, 8.05, ay, 4.4, n, head, body)
        ay += 1.18
    footer(s, foot)
    return s

mockup_slide("05", "Визуализация · Intent Onboarding", "onboarding.jpg", "Один экран — три пути",
    [("1", "Прямой вопрос о цели", "«What are you using … for?» — выбор роли вместо тура по функциям."),
     ("2", "Карточки-роли", "Marketing / HR / Sales — каждая ведёт к своему рабочему пространству."),
     ("3", "Прогресс и «Skip»", "Шаг 3 из 5 + «Skip for now»: контроль и низкий порог входа.")],
    "Референс: intent-driven onboarding")

# =================================================================== 6 HUBSPOT
s = slide()
kicker(s, "Кейс 2 · UI Skins", "06")
title(s, "HubSpot: модульность и upsell через «замочки»")
lead(s, "Единое ядро, разделённое на Hubs. Недоступное не удаляется из интерфейса, а превращается в точку продажи.")
cx = 0.92
cards6 = [("Единое ядро", "Разделение на Marketing, Sales и Service Hubs поверх общей базы контактов и сделок."),
          ("Умное скрытие", "Ненужные функции скрыты или помечены «замочком» — интерфейс остаётся чистым."),
          ("Upsell-тизеры", "Клик по заблокированной функции открывает Demo/Upgrade — disclosure как канал продаж.")]
for h, b in cards6:
    card(s, cx, 2.4, 3.72, 2.2, heading=h, body=b)
    cx += 3.99
takeaway(s, 0.92, 4.95, SW - 1.84, "Вывод для Pitch Avatar:",
         "заблокированная функция — это не «нет», а «ещё не сейчас». «Замочек» упрощает UI и продаёт следующий тариф.")
footer(s)

# =================================================================== 7 MOCKUP modular
mockup_slide("07", "Визуализация · Modular UI & Upsell", "modular.jpg", "Чистое меню + витрина апгрейда",
    [("1", "Активные модули сверху", "Пользователь видит только доступное и нужное прямо сейчас."),
     ("2", "Заблокированное — под «замочком»", "Advanced Analytics, CRM Integrations остаются как тизеры."),
     ("3", "Экран Upgrade to Pro", "Ценность апгрейда — списком фич + явный CTA и «Compare Plans».")],
    "Референс: modular UI & upsell teasers")

# =================================================================== 8 SALESFORCE
s = slide()
kicker(s, "Кейс 3 · UI Skins", "08")
title(s, "Salesforce: «Apps» как настраиваемые UI-скины")
lead(s, "Одно ядро базы данных, поверх которого админ собирает разные приложения под HR, Sales и Marketing — без кода.")
cx = 0.92
cards8 = [("Концепция Apps", "Переключение между HR / Sales / Marketing внутри одного ядра — у каждой роли свой «App»."),
          ("Динамический UI", "Меняется навигация, Home Page и доступные Actions — под задачи роли."),
          ("Lightning App Builder", "Видимость меню и дашбордов настраивается галочками/тумблерами в админке.")]
for h, b in cards8:
    card(s, cx, 2.4, 3.72, 2.2, heading=h, body=b)
    cx += 3.99
takeaway(s, 0.92, 4.95, SW - 1.84, "Вывод для Pitch Avatar:",
         "UI-скин — это конфигурация, а не отдельная сборка. Админ собирает рабочее пространство тумблерами, без разработки.")
footer(s)

# =================================================================== 9 MOCKUP app builder
mockup_slide("09", "Визуализация · Admin App Builder", "appbuilder.jpg", "Конструктор скина для админа",
    [("1", "Theme / Skin", "Enterprise Default, Classic Light, Dark, Custom Skin."),
     ("2", "Sidebar Navigation", "Тумблеры видимости, drag-and-drop порядок, иконки и настройки."),
     ("3", "App Modules", "Включение модулей (Lead Scoring, AI Insights) + Save & Publish.")],
    "Референс: admin app builder")

# =================================================================== 10 COMPARISON TABLE
s = slide()
kicker(s, "Сравнение", "10")
title(s, "Что берём у каждого")
rows = [
    ["Механика", "Monday / ClickUp", "HubSpot", "Salesforce", "→ Pitch Avatar"],
    ["Старт", "Обязательный survey о цели", "Выбор Hub", "Выбор App", "Use Case в БД при регистрации"],
    ["Конфигурация UI", "Шаблоны под ответ", "Модули Hub'ов", "App Builder (no-code)", "Тумблеры модулей (ClickApps)"],
    ["Сложное / платное", "Перевод на Demo", "«Замочки» + upsell", "Права и профили", "«Замочки» с тизером апгрейда"],
    ["Кто настраивает", "Продукт автоматически", "Пользователь / админ", "Админ галочками", "Админ — без разработки"],
    ["Безопасность", "—", "Серверная проверка тарифа", "RLS / профили", "UI-скрытие + 403 на бэкенде"],
]
nrows, ncols = len(rows), len(rows[0])
gt = s.shapes.add_table(nrows, ncols, Inches(0.92), Inches(2.0), Inches(SW - 1.84), Inches(4.6)).table
gt.first_row = False
gt.horz_banding = False
widths = [2.3, 2.35, 2.25, 2.15, 2.45]
for i, wd in enumerate(widths):
    gt.columns[i].width = Inches(wd)
for r in range(nrows):
    gt.rows[r].height = Inches(0.7 if r == 0 else 0.78)
    for c in range(ncols):
        cell = gt.cell(r, c)
        cell.margin_left = Inches(0.14); cell.margin_right = Inches(0.1)
        cell.margin_top = Inches(0.07); cell.margin_bottom = Inches(0.07)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf = cell.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]; run = p.add_run(); run.text = rows[r][c]
        run.font.name = FONT
        if r == 0:
            cell.fill.solid(); cell.fill.fore_color.rgb = WHITE if c < 4 else SOFT
            run.font.size = Pt(10.5); run.font.bold = True
            run.font.color.rgb = (INK if c == 0 else (BRAND if c == 4 else MUTED))
        else:
            run.font.size = Pt(11)
            if c == 0:
                cell.fill.solid(); cell.fill.fore_color.rgb = WHITE
                run.font.bold = True; run.font.color.rgb = INK
            elif c == 4:
                cell.fill.solid(); cell.fill.fore_color.rgb = SOFT
                run.font.color.rgb = NAVYTXT
            else:
                cell.fill.solid(); cell.fill.fore_color.rgb = (ALT if r % 2 == 0 else WHITE)
                run.font.color.rgb = INK2
footer(s)

# =================================================================== 11 GOLDEN RULES
s = slide()
kicker(s, "Принципы", "11")
title(s, "Золотые правила для Pitch Avatar")
rules = [("1", "Субдомен = точка старта", "Use Case записываем в БД при регистрации. URL/лендинг — только триггер, не источник правды."),
         ("2", "Модульность (ClickApps)", "UI-скин — набор тумблеров. Выключили «Лидогенерацию» — она исчезает везде, консистентно."),
         ("3", "Upsell-тизеры", "Скрываем самое сложное, но оставляем «ключи» к апгрейду — «замочки» вместо пустого места."),
         ("4", "Backend Security", "Скрытие в UI всегда подкреплено защитой эндпоинтов (403). UI — это удобство, не безопасность.")]
positions = [(0.92, 2.3), (6.78, 2.3), (0.92, 3.95), (6.78, 3.95)]
for (n, h, b), (x, y) in zip(rules, positions):
    rect(s, x, y, 5.55, 1.45, WHITE, LINE, 1.0, rounded=True, radius=0.06)
    rect(s, x, y, 0.08, 1.45, BRAND, None)
    text(s, x + 0.3, y + 0.22, 0.7, 0.9, [[(n, 26, BRAND, True)]])
    text(s, x + 1.0, y + 0.22, 4.4, 0.4, [[(h, 15, INK, True)]])
    text(s, x + 1.0, y + 0.62, 4.4, 0.7, [[(b, 12, MUTED, False)]], line_spacing=1.12)
takeaway(s, 0.92, 5.65, SW - 1.84, "Главный принцип:",
         "конфигурация важнее кастомных сборок. Одно ядро + декларативные настройки = быстрый запуск любого UI-скина.", h=0.9)
footer(s)

# =================================================================== 12 ROADMAP
s = slide()
kicker(s, "Рекомендации · Roadmap", "12")
title(s, "Как внедрить: три фазы")
phases = [("Фаза 1 · Foundation", "Intent & данные",
           ["Survey на онбординге: 1 вопрос о роли", "Поле use_case в профиле + проброс UTM", "Шаблоны под HR / Sales / Marketing"],
           "Быстрая ценность · минимум зависимостей"),
          ("Фаза 2 · Modular UI", "Скины и «замочки»",
           ["Реестр модулей с тумблерами (flags)", "«Замочки» + экран Upgrade", "Скин по роли поверх ядра"],
           "Зависит от Фазы 1 · растит конверсию"),
          ("Фаза 3 · Admin Builder", "No-code конструктор",
           ["App Builder: видимость пунктов/модулей", "Выбор Theme/Skin, drag-and-drop", "Серверная авторизация (403)"],
           "Для Enterprise · самонастройка")]
cx = 0.92
for kick, h, items, when in phases:
    rect(s, cx, 2.3, 3.72, 3.05, WHITE, LINE, 1.0, rounded=True, radius=0.05)
    text(s, cx + 0.26, 2.52, 3.3, 0.3, [[(kick.upper(), 10.5, BRAND, True)]])
    text(s, cx + 0.26, 2.86, 3.3, 0.4, [[(h, 16, INK, True)]])
    paras = [[("✓  ", 12, GREEN, True), (it, 11.5, INK2, False)] for it in items]
    text(s, cx + 0.26, 3.35, 3.3, 1.3, paras, space_after=7, line_spacing=1.1)
    rect(s, cx + 0.26, 4.92, 3.2, 0.012, LINE, None)
    text(s, cx + 0.26, 5.02, 3.3, 0.3, [[(when, 10.5, MUTED, False)]], line_spacing=1.1)
    cx += 3.99
takeaway(s, 0.92, 5.7, SW - 1.84, "Следующий шаг:",
         "завести эпик «Modular UI & Intent Onboarding» в бэклоге и собрать прототип Фазы 1 в админке pitch-avatar-admin.", h=0.85)
footer(s)

prs.save(OUT)
print("PPTX_OK", OUT, "slides:", len(prs.slides._sldIdLst))
